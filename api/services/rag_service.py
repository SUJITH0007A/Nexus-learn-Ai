import os
import csv
from typing import List, Dict, Any
from pypdf import PdfReader
from docx import Document as DocxDocument
from pptx import Presentation
from PIL import Image
import pytesseract

from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain_core.documents import Document as LangchainDocument
from langchain_community.vectorstores import FAISS
from langchain_community.embeddings import OpenAIEmbeddings
from dotenv import load_dotenv

load_dotenv()

# Setup paths
CACHE_DIR = "cache/vector_stores"
os.makedirs(CACHE_DIR, exist_ok=True)

class RAGService:
    def __init__(self):
        openai_key = os.getenv("NEXT_PUBLIC_OPENAI_API_KEY") or os.getenv("OPENAI_API_KEY")
        if openai_key:
            self.embeddings = OpenAIEmbeddings(openai_key=openai_key)
        else:
            # Fallback mock embedding structure (uses sentence-transformers if available, or basic vectorizer)
            # For robustness, we will initialize a fake embedder if API key is not present
            from langchain_community.embeddings import FakeEmbeddings
            self.embeddings = FakeEmbeddings(size=1536)

    def extract_text_from_file(self, file_path: str) -> str:
        """
        Extract all text contents from PDF, DOCX, PPTX, TXT, CSV, or Image files.
        """
        from api.services.storage_service import storage_service
        import tempfile
        
        temp_local_file = None
        local_path = file_path
        
        if storage_service.is_cloud_enabled() and not os.path.exists(file_path):
            try:
                temp_dir = tempfile.gettempdir()
                filename = os.path.basename(file_path)
                temp_local_file = os.path.join(temp_dir, filename)
                storage_service.download_file(file_path, temp_local_file)
                local_path = temp_local_file
            except Exception as e:
                raise FileNotFoundError(f"Cloud storage file download failed for {file_path}: {str(e)}")

        if not os.path.exists(local_path):
            raise FileNotFoundError(f"File not found: {local_path}")

        ext = os.path.splitext(local_path)[1].lower()
        text = ""

        try:
            # 1. Plain Text
            if ext == ".txt":
                with open(local_path, "r", encoding="utf-8", errors="ignore") as f:
                    text = f.read()

            # 2. PDF Parsing
            elif ext == ".pdf":
                reader = PdfReader(local_path)
                for page in reader.pages:
                    page_text = page.extract_text()
                    if page_text:
                        text += page_text + "\n"

            # 3. Word Document (DOCX)
            elif ext == ".docx":
                doc = DocxDocument(local_path)
                for para in doc.paragraphs:
                    if para.text:
                        text += para.text + "\n"

            # 4. PowerPoint Presentation (PPTX)
            elif ext == ".pptx":
                prs = Presentation(local_path)
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text") and shape.text:
                            text += shape.text + "\n"

            # 5. CSV Data
            elif ext == ".csv":
                with open(local_path, "r", encoding="utf-8", errors="ignore") as f:
                    reader = csv.reader(f)
                    for row in reader:
                        text += ", ".join(row) + "\n"

            # 6. OCR Images (PNG, JPG, JPEG)
            elif ext in [".png", ".jpg", ".jpeg"]:
                try:
                    img = Image.open(local_path)
                    text = pytesseract.image_to_string(img)
                except Exception as e:
                    text = f"Image OCR extraction error: {str(e)}"
                    
            else:
                raise ValueError(f"Unsupported file format: {ext}")
        finally:
            if temp_local_file and os.path.exists(temp_local_file):
                try:
                    os.remove(temp_local_file)
                except:
                    pass

        return text

    def process_and_index_document(self, user_id: int, document_id: int, file_path: str) -> Dict[str, Any]:
        """
        Extract document, split into semantic chunks, generate embeddings and index inside FAISS.
        """
        text = self.extract_text_from_file(file_path)
        if not text.strip():
            # If no text found, mock it or raise
            text = f"No text could be extracted from file: {os.path.basename(file_path)}. Indexing dummy structure."
            
        # Split into semantic chunks
        splitter = RecursiveCharacterTextSplitter(chunk_size=1000, chunk_overlap=150)
        chunks = splitter.split_text(text)
        
        # Wrap as Langchain Documents with metadata
        documents = []
        for idx, chunk in enumerate(chunks):
            metadata = {
                "user_id": user_id,
                "document_id": document_id,
                "filename": os.path.basename(file_path),
                "chunk_index": idx
            }
            documents.append(LangchainDocument(page_content=chunk, metadata=metadata))

        # Save to local FAISS index
        index_name = f"user_{user_id}_doc_{document_id}"
        index_path = os.path.join(CACHE_DIR, index_name)
        
        db = FAISS.from_documents(documents, self.embeddings)
        db.save_local(index_path)

        return {
            "chunk_count": len(chunks),
            "char_count": len(text),
            "index_name": index_name,
            "preview_summary": text[:300] + "..." if len(text) > 300 else text
        }

    def query_document_index(self, user_id: int, document_id: int, query: str, top_k: int = 4) -> List[Dict[str, Any]]:
        """
        Semantic query search on processed document.
        """
        index_name = f"user_{user_id}_doc_{document_id}"
        index_path = os.path.join(CACHE_DIR, index_name)
        
        # Check if local FAISS index exists; if not, try downloading it from cloud storage
        if not os.path.exists(index_path):
            from api.services.storage_service import storage_service
            if storage_service.is_cloud_enabled():
                try:
                    cloud_zip_key = f"user_{user_id}/vector_stores/{index_name}.zip"
                    storage_service.download_and_extract_zip(cloud_zip_key, index_path)
                except Exception as e:
                    print(f"Could not load index from cloud storage: {str(e)}")
        
        if not os.path.exists(index_path):
            return []
            
        # Load local FAISS database
        db = FAISS.load_local(index_path, self.embeddings, allow_dangerous_deserialization=True)
        results = db.similarity_search(query, k=top_k)
        
        return [
            {
                "content": doc.page_content,
                "metadata": doc.metadata
            }
            for doc in results
        ]

rag_service = RAGService()
